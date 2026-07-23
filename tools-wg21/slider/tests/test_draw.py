"""draw.py primitives: import cleanly and pad the box interior correctly."""

from pptx import Presentation
from pptx.util import Inches

import draw
import style


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_draw_box_returns_inner_rect_inset_by_padding():
    cfg = style.load_default()
    box = (5.0, 1.0, 4.0, 2.0)
    inner = draw.draw_box(_slide(), cfg, box, cfg["colors"]["code_bg"], size_pt=14)
    ix, iy, iw, ih = inner
    assert ix > box[0] and iy > box[1]
    assert iw < box[2] and ih < box[3]


def test_draw_box_with_accent_shifts_inner_further_right():
    cfg = style.load_default()
    box = (5.0, 1.0, 4.0, 2.0)
    plain = draw.draw_box(_slide(), cfg, box, cfg["colors"]["callout_bg"], size_pt=18)
    accented = draw.draw_box(_slide(), cfg, box, cfg["colors"]["callout_bg"], size_pt=18,
                             accent_hex=cfg["colors"]["orange"])
    assert accented[0] > plain[0]
