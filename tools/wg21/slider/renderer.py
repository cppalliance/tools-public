"""Slide model -> themed .pptx.

Composition only: this module decides what each slide contains and where each
element sits, delegating every pptx call to draw.py, every number to the style
dict (cfg) and its resolved geometry (geom from layout.py).

Right-panel body flow: each element is its own shape, placed by a running
vertical cursor. Between elements the cursor advances by a gap proportional to
the next element's line height, so spacing scales with font size; when the
element family changes (e.g. paragraph -> list), the boundary gets one extra
line height so consecutive blocks read as separate, not as one run of text.
"""

from __future__ import annotations

import os
import re
import zipfile

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

import draw
import layout
import style
from parser import (
    Bullet,
    Callout,
    CodeBlock,
    ContentSlide,
    Paragraph,
    Subheading,
    TextRun,
    TitleSlide,
    runs_to_plain,
)


def _bold(runs: list[TextRun]) -> list[TextRun]:
    return [TextRun(r.text, True, r.italic, r.code, r.url) for r in runs]


def _italic(runs: list[TextRun]) -> list[TextRun]:
    return [TextRun(r.text, r.bold, True, r.code, r.url) for r in runs]


def _role_font(cfg, role: str) -> str:
    """The font family a chrome role renders in, so fitting measures the right glyphs."""
    f = cfg["fonts"]
    if role == "ts_subtitle":
        return f["italic"]
    if role in ("slide_title", "ts_title"):
        return f["title"]
    return f["label"]


def _fit(cfg, text: str, box, role: str) -> int:
    """Largest size for `role` whose `text` fits `box`, floored at the role's min."""
    floor = cfg["min_sizes"].get(role, cfg["min_sizes"]["fallback"])
    return layout.fit_font_size(cfg, text, box[2], box[3], cfg["sizes"][role], floor,
                                font_family=_role_font(cfg, role))


def _text(slide, cfg, box, runs, size, color, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Place a single-paragraph text box - the chrome around every slide."""
    tf = draw.add_textbox(slide, box, anchor=anchor)
    draw.fill_paragraph(tf.paragraphs[0], runs, cfg, size, color, align=align)


# --- Title / interstitial slide ----------------------------------------------


def render_title_slide(slide, cfg, geom, s: TitleSlide, year: str):
    c = cfg["colors"]
    if s.image_path and os.path.isfile(s.image_path):
        draw.place_image(slide, s.image_path, geom.slide, cfg)
    else:
        draw.add_rect(slide, geom.slide, c["background"])

    _text(slide, cfg, geom.ts_title, s.title, _fit(cfg, runs_to_plain(s.title), geom.ts_title, "ts_title"), c["white"])
    if s.subtitle:
        sub = [TextRun(s.subtitle, italic=True)]
        _text(slide, cfg, geom.ts_subtitle, sub, _fit(cfg, s.subtitle, geom.ts_subtitle, "ts_subtitle"), c["orange"])
    _text(slide, cfg, geom.ts_year, [TextRun(year)], cfg["sizes"]["year"], c["muted"], anchor=MSO_ANCHOR.BOTTOM)


# --- Content slide ------------------------------------------------------------


def render_content_slide(slide, cfg, geom, s: ContentSlide, page_no: int, total: int):
    c = cfg["colors"]
    # Fill the whole slide first, then lay the navy left panel on top - no seam.
    draw.add_rect(slide, geom.slide, c["background"])
    draw.add_rect(slide, geom.left_panel, c["navy"])
    draw.place_image(slide, s.image_path, geom.image, cfg)

    _text(slide, cfg, geom.label, [TextRun(s.section_label, bold=True)],
          _fit(cfg, s.section_label, geom.label, "section_label"), c["orange"])
    # Title lives on the right panel now: large, centered, above the body.
    _text(slide, cfg, geom.title, s.title,
          _fit(cfg, runs_to_plain(s.title), geom.title, "slide_title"), c["white"],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    _render_body(slide, cfg, geom, s.body)

    _text(slide, cfg, geom.page_number, [TextRun(f"{page_no} / {total}")],
          cfg["sizes"]["page_number"], c["muted"], anchor=MSO_ANCHOR.BOTTOM, align=PP_ALIGN.RIGHT)


# --- Right-panel body flow ----------------------------------------------------


def _role(el) -> str:
    if isinstance(el, CodeBlock):
        return "code"
    if isinstance(el, Subheading):
        return "subheading"
    if isinstance(el, Callout):
        return "callout"
    if isinstance(el, Bullet):
        return "subbullet" if el.level > 0 else ("numbered" if el.ordered else "bullet")
    return "body"


def _plain(el) -> str:
    return el.text if isinstance(el, CodeBlock) else runs_to_plain(el.runs)


def _role_size(cfg, role: str, scale: float) -> int:
    floor = cfg["min_sizes"].get(role, cfg["min_sizes"]["fallback"])
    return max(floor, int(round(cfg["sizes"][role] * scale)))


def _indent(cfg, el, size: int) -> float:
    if isinstance(el, Bullet):
        return style.em(size, cfg["body"]["bullet_indent"] + el.level * cfg["body"]["bullet_indent_step"])
    return 0.0


def _text_width(cfg, el, size: int, content_w: float) -> float:
    if isinstance(el, CodeBlock):
        return content_w - 2 * style.em(size, cfg["body"]["box_pad_x"])
    if isinstance(el, Callout):
        return content_w - 2 * style.em(size, cfg["body"]["box_pad_x"]) - style.em(size, cfg["body"]["accent_bar"])
    return content_w - _indent(cfg, el, size)


def _elem_font(cfg, el) -> str:
    """The font family a body element renders in, for accurate height measurement."""
    f = cfg["fonts"]
    if isinstance(el, CodeBlock):
        return f["mono"]
    if isinstance(el, Callout):
        return f["italic"]
    return f["body"]


def _elem_height(cfg, el, size: int, content_w: float) -> float:
    h = layout.text_height_in(cfg, _plain(el), _text_width(cfg, el, size, content_w),
                              size, mono=isinstance(el, CodeBlock), font_family=_elem_font(cfg, el))
    if isinstance(el, (Callout, CodeBlock)):
        h += 2 * style.em(size, cfg["body"]["box_pad_y"])
    return h


def _gap(cfg, size: int) -> float:
    return cfg["text"]["gap_ratio"] * layout.line_height_in(cfg, size)


def _family(el) -> str:
    """Element family for inter-block spacing. All list items are one family, so
    bullet-to-subbullet and bullet-to-numbered boundaries stay tight; paragraph,
    subheading, callout, and code are each their own family."""
    if isinstance(el, Bullet):
        return "list"
    if isinstance(el, Subheading):
        return "subheading"
    if isinstance(el, Callout):
        return "callout"
    if isinstance(el, CodeBlock):
        return "code"
    return "body"


def _lead_gap(cfg, prev_el, el, size: int) -> float:
    """Vertical gap above `el`, given the element before it. Zero for the first
    element; the base gap between same-family neighbours; the base gap plus one
    line height of `el` when it opens a new element family, so a block boundary
    reads as a break rather than another line of the previous block."""
    if prev_el is None:
        return 0.0
    gap = _gap(cfg, size)
    if _family(prev_el) != _family(el):
        gap += layout.line_height_in(cfg, size)
    return gap


def _body_scale(cfg, body, content_w: float, content_h: float) -> float:
    # Heights and gaps (base and family-transition) are all linear in font size,
    # so total height scales linearly; one ratio brings an overset body back
    # inside the panel.
    total = 0.0
    prev = None
    for el in body:
        size = _role_size(cfg, _role(el), 1.0)
        total += _lead_gap(cfg, prev, el, size)
        total += _elem_height(cfg, el, size, content_w)
        prev = el
    if total <= content_h or total <= 0:
        return 1.0
    return max(cfg["text"]["min_body_scale"], content_h / total)


def _marker(cfg, el: Bullet) -> str:
    if el.ordered and el.number is not None:
        return f"{el.number}."
    return cfg["markers"]["circle"] if el.level > 0 else cfg["markers"]["square"]


def _render_body(slide, cfg, geom, body):
    c = cfg["colors"]
    cx, cy, cw, ch = geom.content
    scale = _body_scale(cfg, body, cw, ch)
    y = cy
    prev = None
    for el in body:
        role = _role(el)
        size = _role_size(cfg, role, scale)
        y += _lead_gap(cfg, prev, el, size)
        h = _elem_height(cfg, el, size, cw)

        if isinstance(el, CodeBlock):
            inner = draw.draw_box(slide, cfg, (cx, y, cw, h), c["code_bg"], size)
            tf = draw.add_textbox(slide, inner)
            for j, line in enumerate(el.text.split("\n")):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                draw.fill_paragraph(p, [TextRun(line, code=True)], cfg, size, c["white"], code_color=c["code_text"])

        elif isinstance(el, Callout):
            inner = draw.draw_box(slide, cfg, (cx, y, cw, h), c["callout_bg"], size, accent_hex=c["orange"])
            tf = draw.add_textbox(slide, inner)
            draw.fill_paragraph(tf.paragraphs[0], _italic(el.runs), cfg, size, c["white"], mono_bg=c["code_bg"])

        elif isinstance(el, Subheading):
            tf = draw.add_textbox(slide, (cx, y, cw, h))
            draw.fill_paragraph(tf.paragraphs[0], _bold(el.runs), cfg, size, c["white"])

        elif isinstance(el, Bullet):
            indent = _indent(cfg, el, size)
            marker_w = style.em(size, cfg["body"]["marker_col"])
            mtf = draw.add_textbox(slide, (cx + indent - marker_w, y, marker_w, h))
            draw.fill_paragraph(mtf.paragraphs[0], [TextRun(_marker(cfg, el), bold=True)], cfg, size, c["orange"])
            ttf = draw.add_textbox(slide, (cx + indent, y, cw - indent, h))
            runs = _bold(el.runs) if el.highlight else el.runs
            draw.fill_paragraph(ttf.paragraphs[0], runs, cfg, size, c["orange"], mono_bg=c["code_bg"])

        else:  # Paragraph
            tf = draw.add_textbox(slide, (cx, y, cw, h))
            draw.fill_paragraph(tf.paragraphs[0], el.runs, cfg, size, c["white"], mono_bg=c["code_bg"])

        y += h
        prev = el


# --- Speaker notes ------------------------------------------------------------


def _apply_notes(slide, notes: list[str]) -> None:
    # Touching slide.notes_slide creates it, so a slide with no notes never
    # gets a notes page.
    if not notes:
        return
    tf = slide.notes_slide.notes_text_frame
    for i, para in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para


# --- Deck ---------------------------------------------------------------------


def build_presentation(slides, year: str, cfg: dict | None = None) -> Presentation:
    cfg = cfg or style.load_default()
    geom = layout.geometry(cfg)
    prs = Presentation()
    prs.slide_width = Inches(cfg["slide"]["width"])
    prs.slide_height = Inches(cfg["slide"]["height"])
    blank = prs.slide_layouts[6]

    total = len(slides)
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        if isinstance(s, TitleSlide):
            render_title_slide(slide, cfg, geom, s, year)
        else:
            render_content_slide(slide, cfg, geom, s, i + 1, total)
        _apply_notes(slide, getattr(s, "notes", []))
    return prs


def _patch_theme_hyperlinks(path: str, hex_color: str) -> None:
    # PowerPoint and LibreOffice color hyperlinked runs from the theme's
    # hlink/folHlink slots, not the run fill, so links would show theme blue
    # unless we rewrite those slots in the saved file.
    try:
        tmp = path + ".tmp"
        pat = re.compile(r"<a:(hlink|folHlink)>.*?</a:\1>", re.DOTALL)

        def repl(m):
            tag = m.group(1)
            return f'<a:{tag}><a:srgbClr val="{hex_color}"/></a:{tag}>'

        with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.startswith("ppt/theme/") and item.filename.endswith(".xml"):
                    data = pat.sub(repl, data.decode("utf-8")).encode("utf-8")
                zout.writestr(item, data)
        os.replace(tmp, path)
    except Exception:
        if os.path.exists(path + ".tmp"):
            os.remove(path + ".tmp")


def save_presentation(prs: Presentation, path: str, accent_hex: str | None = None) -> None:
    prs.save(path)
    _patch_theme_hyperlinks(path, accent_hex or style.load_default()["colors"]["orange"])
