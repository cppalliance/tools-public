"""pptx drawing primitives: rectangles, text boxes, run styling, boxes, images.

This module knows how to put shapes on a python-pptx slide from a style dict.
It does not know what a deck contains - that is renderer.py's job.
"""

from __future__ import annotations

import io
import os

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from parser import TextRun
from style import em

# rPr child order in the OOXML schema; <a:highlight> must precede these.
_RPR_TAIL = ("a:latin", "a:ea", "a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _jpeg_quality(level: int) -> int:
    # Map the style's 1-9 compression scale to a JPEG quality percentage:
    # 1 = maximum compression (~10), 9 = minimal compression (~95). Level 0
    # disables conversion upstream, so it never reaches here.
    return round(10 + (level - 1) * 10.625)


def _png_to_jpeg_stream(path: str, cfg: dict) -> io.BytesIO:
    # Re-encode a PNG as an in-memory JPEG, flattening any transparency onto the
    # style's slide background color. The file on disk is never modified.
    with Image.open(path) as img:
        if img.mode in ("RGBA", "LA", "P"):
            img = img.convert("RGBA")
            flat = Image.new("RGB", img.size, _hex_to_rgb(cfg["colors"]["background"]))
            flat.paste(img, mask=img.split()[-1])
            img = flat
        else:
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=_jpeg_quality(cfg["images"]["jpeg_quality"]))
    buf.seek(0)
    return buf


def add_rect(slide, box, fill_hex):
    x, y, w, h = box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def add_textbox(slide, box, anchor=MSO_ANCHOR.TOP, wrap=True):
    x, y, w, h = box
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def set_highlight(run, hex_color: str) -> None:
    # python-pptx has no highlight API, so inject <a:highlight> at the schema
    # position by hand. Any failure degrades to plain colored text.
    try:
        rpr = run._r.get_or_add_rPr()
        for old in rpr.findall(qn("a:highlight")):
            rpr.remove(old)
        hl = rpr.makeelement(qn("a:highlight"), {})
        hl.append(rpr.makeelement(qn("a:srgbClr"), {"val": hex_color}))
        anchor = next((c for c in rpr if c.tag in (qn(t) for t in _RPR_TAIL)), None)
        if anchor is not None:
            anchor.addprevious(hl)
        else:
            rpr.append(hl)
    except Exception:
        pass


def apply_run(run, tr: TextRun, cfg: dict, default_color: str, mono_bg: str | None, code_color: str) -> None:
    c = cfg["colors"]
    f = run.font
    if tr.code:
        f.name = cfg["fonts"]["mono"]
    elif tr.italic:
        f.name = cfg["fonts"]["italic"]
    else:
        f.name = cfg["fonts"]["body"]
    f.bold = tr.bold
    f.italic = tr.italic
    if tr.url:
        f.color.rgb = rgb(c["orange"])
        f.underline = True
        try:
            run.hyperlink.address = tr.url
        except Exception:
            pass
    elif tr.code:
        f.color.rgb = rgb(code_color)
    elif tr.italic and default_color == c["white"]:
        # Dim only body italics; a caller-set color (orange subtitle, links) stays.
        f.color.rgb = rgb(c["italic_dim"])
    else:
        f.color.rgb = rgb(default_color)
    if tr.code and mono_bg:
        set_highlight(run, mono_bg)


def fill_paragraph(p, runs, cfg, size_pt, default_color, mono_bg=None, align=PP_ALIGN.LEFT, code_color=None):
    p.alignment = align
    p.line_spacing = cfg["text"]["line_spacing"]
    code_color = code_color or cfg["colors"]["orange"]
    for tr in runs or [TextRun("")]:
        # A pptx run cannot hold a newline, so split hard breaks into runs.
        for i, part in enumerate(tr.text.split("\n")):
            if i > 0:
                p.add_line_break()
            r = p.add_run()
            r.text = part
            r.font.size = Pt(size_pt)
            apply_run(r, TextRun(part, tr.bold, tr.italic, tr.code, tr.url), cfg, default_color, mono_bg, code_color)


def draw_box(slide, cfg, box, fill_hex, size_pt, accent_hex=None):
    """Fill a background box, draw an optional left accent bar, return the padded inner box."""
    x, y, w, h = box
    add_rect(slide, box, fill_hex)
    pad_x = em(size_pt, cfg["body"]["box_pad_x"])
    pad_y = em(size_pt, cfg["body"]["box_pad_y"])
    inner_x = x + pad_x
    if accent_hex is not None:
        bar = em(size_pt, cfg["body"]["accent_bar"])
        add_rect(slide, (x, y, bar, h), accent_hex)
        inner_x = x + bar + pad_x
    return (inner_x, y + pad_y, w - (inner_x - x) - pad_x, h - 2 * pad_y)


def place_image(slide, path, box, cfg=None):
    """Draw the image if the file exists; otherwise leave the background showing.

    When a style is supplied with images.jpeg_quality > 0 and the source is a
    PNG, the image is re-encoded to JPEG in memory before embedding - shrinking
    the .pptx while leaving the original PNG on disk untouched. Any other format,
    or quality 0, embeds the file as-is. A conversion failure falls back to the
    original PNG rather than dropping the image."""
    if not (path and os.path.isfile(path)):
        return
    try:
        x, y, w, h = box
        source = path
        if cfg and cfg.get("images", {}).get("jpeg_quality", 0) > 0 and path.lower().endswith(".png"):
            try:
                source = _png_to_jpeg_stream(path, cfg)
            except Exception:
                source = path
        slide.shapes.add_picture(source, Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception:
        pass
