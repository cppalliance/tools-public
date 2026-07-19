"""Tests for speaker-notes support (the second `---` section).

Parser tests assert notes land on the slide model at the right boundaries;
renderer tests build a real `.pptx`, reopen it, and assert the notes pane text.
A regression test proves single-`---` decks are untouched.
"""

import os

from pptx import Presentation

from parser import Bullet, ContentSlide, Paragraph, TitleSlide, parse_markdown
from renderer import build_presentation, save_presentation

_HERE = os.path.dirname(os.path.abspath(__file__))
GOLDEN = os.path.join(_HERE, "golden", "notes.md")
TEST_MD = os.path.join(os.path.dirname(_HERE), "test.md")


def _read(path: str) -> str:
    with open(path, encoding="utf-8") as fh:
        return fh.read()


def _golden_slides():
    return parse_markdown(_read(GOLDEN), base_dir=os.path.dirname(GOLDEN))


# --- Parser: structure --------------------------------------------------------


def test_parses_expected_slide_sequence():
    slides = _golden_slides()
    assert len(slides) == 6
    kinds = [type(s).__name__ for s in slides]
    assert kinds == [
        "TitleSlide",
        "ContentSlide",
        "ContentSlide",
        "ContentSlide",
        "TitleSlide",
        "ContentSlide",
    ]


# --- Parser: title slides -----------------------------------------------------


def test_title_slide_with_notes():
    cover = _golden_slides()[0]
    assert isinstance(cover, TitleSlide)
    assert cover.subtitle == "The subtitle line."
    assert cover.notes == ["Cover speaker notes paragraph one. Second sentence, same paragraph."]


def test_title_slide_without_notes():
    interstitial = _golden_slides()[4]
    assert isinstance(interstitial, TitleSlide)
    assert interstitial.subtitle == "A divider with no notes."
    assert interstitial.notes == []


# --- Parser: content slides ---------------------------------------------------


def test_content_slide_with_body_and_notes():
    first = _golden_slides()[1]
    assert isinstance(first, ContentSlide)
    assert first.section_label == "FIRST SECTION"
    assert any(isinstance(b, Paragraph) for b in first.body)
    assert any(isinstance(b, Bullet) for b in first.body)
    assert first.notes == ["Notes for the first content slide."]


def test_content_slide_body_only_has_no_notes():
    second = _golden_slides()[2]
    assert isinstance(second, ContentSlide)
    assert second.body  # body present
    assert second.notes == []


def test_content_slide_rich_notes_flatten_to_text():
    third = _golden_slides()[3]
    assert third.notes == [
        "First notes paragraph with bold and italic.",
        "Second notes paragraph.",
        "- Note bullet one",
        "- Note bullet two",
    ]


def test_notes_run_to_end_of_file_without_swallowing_body():
    last = _golden_slides()[5]
    assert last.section_label == "FOURTH SECTION"
    assert any(isinstance(b, Paragraph) for b in last.body)
    assert last.notes == ["These notes run to the end of the file. Second line of final notes."]


# --- Renderer: notes land in the .pptx notes pane -----------------------------


def test_renderer_writes_notes_to_pptx(tmp_path):
    slides = _golden_slides()
    prs = build_presentation(slides, year="2026")
    out = str(tmp_path / "notes.pptx")
    save_presentation(prs, out)

    reopened = list(Presentation(out).slides)

    assert reopened[0].has_notes_slide
    assert (
        reopened[0].notes_slide.notes_text_frame.text
        == "Cover speaker notes paragraph one. Second sentence, same paragraph."
    )

    assert reopened[1].notes_slide.notes_text_frame.text == "Notes for the first content slide."

    # Body-only and interstitial slides get no notes page at all.
    assert not reopened[2].has_notes_slide
    assert not reopened[4].has_notes_slide

    assert reopened[3].notes_slide.notes_text_frame.text == "\n".join(
        [
            "First notes paragraph with bold and italic.",
            "Second notes paragraph.",
            "- Note bullet one",
            "- Note bullet two",
        ]
    )

    assert (
        reopened[5].notes_slide.notes_text_frame.text
        == "These notes run to the end of the file. Second line of final notes."
    )


# --- Regression: single-`---` decks are untouched -----------------------------


def test_existing_deck_has_no_notes(tmp_path):
    slides = parse_markdown(_read(TEST_MD), base_dir=os.path.dirname(TEST_MD))
    assert len(slides) == 5
    assert all(not getattr(s, "notes", []) for s in slides)

    prs = build_presentation(slides, year="2026")
    out = str(tmp_path / "test.pptx")
    save_presentation(prs, out)

    assert all(not sl.has_notes_slide for sl in Presentation(out).slides)
