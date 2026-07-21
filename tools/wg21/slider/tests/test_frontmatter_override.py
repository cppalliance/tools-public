"""A deck's frontmatter theme: override reaches the rendered .pptx."""

import os

from pptx import Presentation

import style
from parser import parse_markdown
from renderer import build_presentation, save_presentation

_DECK = """---
theme:
  inherits: default
  colors:
    orange: "D64500"
---

## ONLY SECTION

The title

---

Body paragraph.
"""


def test_frontmatter_orange_override_colors_the_section_label(tmp_path):
    fm, body = style.extract_frontmatter(_DECK)
    cfg = style.resolve(None, fm.get("theme"))
    assert cfg["colors"]["orange"] == "D64500"

    slides = parse_markdown(body)
    out = str(tmp_path / "d.pptx")
    save_presentation(build_presentation(slides, year="2026", cfg=cfg), out, accent_hex=cfg["colors"]["orange"])

    prs = Presentation(out)
    seen = set()
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.type is not None:
                    seen.add(str(run.font.color.rgb))
    assert "D64500" in seen
